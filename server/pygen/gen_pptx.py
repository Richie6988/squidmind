#!/usr/bin/env python3
"""gen_pptx.py <spec.json> — generate a deck, template-inherited or blank.

Spec: { "template": "/abs/path.pptx" | null, "output": "/abs/out.pptx",
        "title": str?, "author": str?, "asset_root": "/abs/aquarium"?,
        "slides": [ { "title": str, "bullets": [str]?, "body": str?,
                      "notes": str?, "layout": "title"|"content"?,
                      "image": str | {"path": str, "caption": str?}?,
                      "table": {"headers": [str], "rows": [[str]]}?,
                      "chart": {"type": "bar"|"line"|"pie",
                                "categories": [str],
                                "series": [{"name": str, "values": [num]}]}? } ] }

Rich content (image/table/chart) is placed in the content zone under the
title; when bullets/body coexist with rich content, text takes the upper
half and the rich element the lower half. Image paths resolve against
asset_root (aquarium) and are confined to it.

Strategy — maximum inheritance, graceful degradation:
  1. Open the template with python-pptx: masters, layouts, theme colors,
     fonts, logos all come along for FREE.
  2. Drop the template's existing slides (they were style reference, not
     content) via the sldIdLst XML — the documented python-pptx recipe.
  3. For each spec slide, pick the best layout by PLACEHOLDER SHAPES:
     a layout with a TITLE + BODY placeholder for content slides, a
     CENTER_TITLE (or first) layout for title slides. Fill placeholders so
     the template's own positioning/typography applies. If a needed
     placeholder is missing, fall back to a plain textbox.
Result JSON on stdout: {"ok": true, "outputPath": ..., "slides": N}
"""
import json, sys, copy

def fail(msg):
    print(json.dumps({"ok": False, "error": str(msg)})); sys.exit(0)

import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
except Exception as e:
    fail(f"python-pptx not available: {e}")

ASSET_ROOT = None

def resolve_asset(p):
    """Absolute-or-aquarium-relative path, confined to ASSET_ROOT."""
    if not p: raise ValueError("empty image path")
    cand = p if os.path.isabs(p) else os.path.join(ASSET_ROOT or "", p)
    real = os.path.realpath(cand)
    if ASSET_ROOT and not real.startswith(os.path.realpath(ASSET_ROOT) + os.sep):
        raise ValueError(f"image path escapes the aquarium: {p}")
    if not os.path.isfile(real):
        raise ValueError(f"image not found: {p}")
    return real

PH_TITLE  = {13, 1, 3}   # TITLE(13 in some versions?) — use idx-agnostic type names below
try:
    from pptx.enum.shapes import PP_PLACEHOLDER
except Exception:
    PP_PLACEHOLDER = None

def ph_types(layout):
    out = {}
    for ph in layout.placeholders:
        out[ph.placeholder_format.type] = ph.placeholder_format.idx
    return out

def is_title_type(t):
    if PP_PLACEHOLDER is None: return False
    return t in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)

def is_body_type(t):
    if PP_PLACEHOLDER is None: return False
    return t in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)

def pick_layouts(prs):
    title_layout, content_layout = None, None
    for layout in prs.slide_layouts:
        types = [ph.placeholder_format.type for ph in layout.placeholders]
        has_title = any(is_title_type(t) for t in types)
        has_body  = any(is_body_type(t) for t in types)
        if has_title and has_body and content_layout is None:
            content_layout = layout
        if has_title and not has_body and title_layout is None:
            title_layout = layout
    # Fallbacks: conventional indices 0 (title) / 1 (title+content), else any
    layouts = list(prs.slide_layouts)
    if title_layout is None:
        title_layout = layouts[0] if layouts else None
    if content_layout is None:
        content_layout = layouts[1] if len(layouts) > 1 else title_layout
    return title_layout, content_layout

def drop_existing_slides(prs):
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        xml_slides.remove(sld)

def set_title(slide, text):
    for ph in slide.placeholders:
        if is_title_type(ph.placeholder_format.type):
            ph.text = text or ""
            return True
    return False

def set_body(slide, bullets, body):
    for ph in slide.placeholders:
        if is_body_type(ph.placeholder_format.type):
            tf = ph.text_frame
            items = bullets if bullets else ([body] if body else [])
            if not items:
                tf.text = ""
                return True
            tf.text = str(items[0])
            for it in items[1:]:
                p = tf.add_paragraph()
                p.text = str(it)
            return True
    return False

def content_zone(prs, upper_used):
    """(left, top, width, height) EMU for the rich element."""
    W, H = prs.slide_width, prs.slide_height
    left  = Emu(int(W * 0.06))
    width = Emu(int(W * 0.88))
    top_f = 0.52 if upper_used else 0.24
    top    = Emu(int(H * top_f))
    height = Emu(int(H * (0.92 - top_f)))
    return left, top, width, height

def add_image(prs, slide, spec, upper_used):
    path = spec if isinstance(spec, str) else spec.get("path")
    cap  = None if isinstance(spec, str) else spec.get("caption")
    real = resolve_asset(path)
    left, top, width, height = content_zone(prs, upper_used)
    # Fit by width, cap by zone height (python-pptx keeps ratio w/ one dim,
    # so probe the native ratio first)
    from pptx.util import Emu as _E
    pic = slide.shapes.add_picture(real, left, top, width=width)
    if pic.height > height:
        ratio = height / pic.height
        pic.height = int(pic.height * ratio)
        pic.width  = int(pic.width * ratio)
        pic.left   = left + int((width - pic.width) / 2)
    if cap:
        cb = slide.shapes.add_textbox(left, top + pic.height + _E(50000), width, _E(300000))
        tf = cb.text_frame; tf.text = str(cap)
        for p in tf.paragraphs:
            for r in p.runs: r.font.size = Pt(11); r.font.italic = True

def add_table(prs, slide, spec, upper_used):
    headers = spec.get("headers") or []
    rows    = spec.get("rows") or []
    if not headers and rows: headers = ["" for _ in rows[0]]
    nrows, ncols = len(rows) + (1 if headers else 0), max(1, len(headers) or (len(rows[0]) if rows else 1))
    left, top, width, height = content_zone(prs, upper_used)
    shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = shape.table
    r0 = 0
    if headers:
        for c, h in enumerate(headers[:ncols]):
            cell = tbl.cell(0, c); cell.text = str(h)
            for p in cell.text_frame.paragraphs:
                for run in p.runs: run.font.bold = True; run.font.size = Pt(12)
        r0 = 1
    for ri, row in enumerate(rows):
        for c, v in enumerate(row[:ncols]):
            cell = tbl.cell(ri + r0, c); cell.text = str(v)
            for p in cell.text_frame.paragraphs:
                for run in p.runs: run.font.size = Pt(11)

CHART_TYPES = None
def add_chart(prs, slide, spec, upper_used):
    global CHART_TYPES
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    if CHART_TYPES is None:
        CHART_TYPES = {"bar": XL_CHART_TYPE.COLUMN_CLUSTERED, "line": XL_CHART_TYPE.LINE_MARKERS, "pie": XL_CHART_TYPE.PIE}
    ct = CHART_TYPES.get(str(spec.get("type", "bar")).lower(), CHART_TYPES["bar"])
    data = CategoryChartData()
    data.categories = [str(c) for c in (spec.get("categories") or [])]
    series = spec.get("series") or []
    for s in series:
        data.add_series(str(s.get("name", "")), [float(v) for v in (s.get("values") or [])])
    left, top, width, height = content_zone(prs, upper_used)
    gframe = slide.shapes.add_chart(ct, left, top, width, height, data)
    chart = gframe.chart
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.include_in_layout = False

def add_textbox(slide, text, top_in, size, bold=False):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top_in), Inches(12.0), Inches(1.0))
    tf = box.text_frame
    tf.text = text or ""
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
    return box

def main():
    global ASSET_ROOT
    spec = json.load(open(sys.argv[1], encoding="utf-8"))
    ASSET_ROOT = spec.get("asset_root")
    prs = Presentation(spec["template"]) if spec.get("template") else Presentation()
    if spec.get("title"):  prs.core_properties.title  = spec["title"]
    if spec.get("author"): prs.core_properties.author = spec["author"]
    drop_existing_slides(prs)
    title_layout, content_layout = pick_layouts(prs)
    if title_layout is None:
        fail("template has no usable slide layouts")

    n = 0
    for s in spec.get("slides", []):
        is_title = s.get("layout") == "title" or (not s.get("bullets") and not s.get("body")
                    and not s.get("image") and not s.get("table") and not s.get("chart"))
        layout = title_layout if is_title else content_layout
        slide = prs.slides.add_slide(layout)
        if not set_title(slide, s.get("title", "")):
            add_textbox(slide, s.get("title", ""), 0.4, 30, bold=True)
        has_rich = bool(s.get("image") or s.get("table") or s.get("chart"))
        has_text = bool(s.get("bullets") or s.get("body"))
        if is_title:
            if s.get("body") and not set_body(slide, None, s.get("body")):
                add_textbox(slide, s.get("body"), 4.2, 18)
        elif has_text:
            if not set_body(slide, s.get("bullets"), s.get("body")):
                text = "\n".join(s.get("bullets") or ([s.get("body")] if s.get("body") else []))
                add_textbox(slide, text, 1.6, 16)
        # Rich content — lower half when text is above, full zone otherwise
        if s.get("image"): add_image(prs, slide, s["image"], has_text)
        if s.get("table"): add_table(prs, slide, s["table"], has_text or bool(s.get("image")))
        if s.get("chart"): add_chart(prs, slide, s["chart"], has_text or bool(s.get("image")) or bool(s.get("table")))
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = str(s["notes"])
        n += 1

    prs.save(spec["output"])
    print(json.dumps({"ok": True, "outputPath": spec["output"], "slides": n, "engine": "python-pptx (template-inherited)" if spec.get("template") else "python-pptx (blank)"}))

if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        fail(e)
